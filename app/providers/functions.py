from pathlib import Path
import requests
from reportlab.lib.pagesizes import A4, landscape, mm
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from secrets import token_urlsafe
from app.models.tables import Worksheet_Content
from pptx import Presentation
from pptx.util import Mm
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from PIL import UnidentifiedImageError


ALLOWED_EXTENSIONS = {'xlsx', 'xls'}
def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def image_id_generator():
    image_id = token_urlsafe(40)
    image_id_db = Worksheet_Content.query.filter(Worksheet_Content.image_id == image_id).first()
    if image_id_db:
        image_id_generator()
    return image_id

def pdf_generator(capa, content, image_id):
    #try:
    colunas = content['colunas']
    linhas = content['conteudo']

    endereco_column = None
    latitude_column = ''
    longitude_column = ''
    codigo_column = None
    foto_column = None
    other_columns = []

    for coluna in colunas:
        if coluna.lower().rsplit(' ')[0] == 'endereco' or coluna.lower().rsplit(' ')[0] == 'endereço' or coluna.lower().rsplit(' ')[0] == 'direccion' or coluna.lower().rsplit(' ')[0] == 'dirección' or coluna.lower().rsplit(' ')[0] == 'address':
            endereco_column = coluna
        elif coluna.lower().rsplit(' ')[0] == 'latitude' or coluna.lower().rsplit(' ')[0] == 'latitud':
            latitude_column = coluna
        elif coluna.lower().rsplit(' ')[0] == 'longitude' or coluna.lower().rsplit(' ')[0] == 'longitud':
            longitude_column = coluna
        elif coluna.lower().rsplit(' ')[0] == 'cod.' or coluna.lower().rsplit(' ')[0] == 'codigo' or coluna.lower().rsplit(' ')[0] == 'código' or coluna.lower().rsplit(' ')[0] == 'code' or coluna.lower().rsplit(' ')[0] == 'cod' or coluna.lower().rsplit(' ')[0] == 'cód' or coluna.lower().rsplit(' ')[0] == 'cód.':
            codigo_column = coluna
        elif coluna.lower().rsplit(' ')[0] == 'foto' or coluna.lower().rsplit(' ')[0] == 'imagem' or coluna.lower().rsplit(' ')[0] == 'imagen' or coluna.lower().rsplit(' ')[0] == 'fotografia' or coluna.lower().rsplit(' ')[0] == 'fotografía' or coluna.lower().rsplit(' ')[0] == 'image' or coluna.lower().rsplit(' ')[0] == 'picture' or coluna.lower().rsplit(' ')[0] == 'photo':
            foto_column = coluna
        else:
            other_columns.append(coluna)
    if not endereco_column:
        return False, 'A coluna do endereço não foi reconhecida. A planilha deve fornecer uma coluna de nome "Endereco", "Endereço", "Direccion", "Dirección" ou "Address".'
    if not latitude_column:
        return False, 'A coluna da latitude não foi reconhecida. A planilha deve fornecer uma coluna de nome "Latitude" ou "Latitud".'
    if not longitude_column:
        return False, 'A coluna da longitude não foi reconhecida. A planilha deve fornecer uma coluna de nome "Longitude" ou "Longitud".'
    if not codigo_column:
        return False, 'A coluna do código não foi reconhecida. A planilha deve fornecer uma coluna de nome "Cod.", "Cód.", "Cod", "Cód", "Codigo", "Código" ou "Code".'
    if not foto_column:
        return False, 'A coluna da foto não foi reconhecida. A planilha deve fornecer uma coluna de nome "Foto", "Imagem", "Image", "Picture", "Photo", "Imagen" ou "Fotografia".'

    # Gerando PDF
    pdf = canvas.Canvas(f'app/static/media/pdf/{image_id}.pdf', (400*mm, 220*mm))
    # Gerando PPTX
    apresentacao = Presentation()
    apresentacao.slide_height = Mm(220)
    apresentacao.slide_width = Mm(400)

    # Capa
    
    if not capa['nome']:
        return False, 'Você deve fornecer um nome para o book.'
    else:
        # PDF
        pdf.drawImage('app/static/media/pdf_provider_images/capa_template.jpg', 0, 0, 400*mm, 220*mm)
        pdf.setFont('Helvetica', 10*mm)
        
        # Inicializando slide da capa
        slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
        capa_template = slide.shapes.add_picture('app/static/media/pdf_provider_images/capa_template.jpg', Mm(0), Mm(0), height=Mm(220), width=Mm(400))
        
        if not capa['cliente'] and not capa['pessoa']:
            pdf.drawString(105*mm,25*mm, capa['nome'])

            capa_nome = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
            capa_nome_text_frame = capa_nome.text_frame
            capa_nome_text_frame.clear()
            capa_nome_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            capa_nome_text = capa_nome_text_frame.paragraphs[0].add_run()
            capa_nome_text.text = capa['nome']
            capa_nome_text.font.name = 'Helvetica'
            capa_nome_text.font.size = Mm(10)
            
        else:
            if capa['cliente'] and capa['pessoa']:
                texto = f'A\C {capa["pessoa"]} - {capa["cliente"]}'
                if len(texto) > 48:
                    pdf.drawString(105*mm,55*mm, capa['nome'])
                    
                    capa_nome = slide.shapes.add_textbox(Mm(102), Mm(155), Mm(200), Mm(10))
                    capa_nome_text_frame = capa_nome.text_frame
                    capa_nome_text_frame.clear()
                    capa_nome_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    capa_nome_text = capa_nome_text_frame.paragraphs[0].add_run()
                    capa_nome_text.text = capa['nome']
                    capa_nome_text.font.name = 'Helvetica'
                    capa_nome_text.font.size = Mm(10)

                    pdf.drawString(105*mm,40*mm, texto[0:48])

                    texto_pptx1 = slide.shapes.add_textbox(Mm(102), Mm(170), Mm(200), Mm(10))
                    texto_pptx1_text_frame = texto_pptx1.text_frame
                    texto_pptx1_text_frame.clear()
                    texto_pptx1_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    texto_pptx1_text = texto_pptx1_text_frame.paragraphs[0].add_run()
                    texto_pptx1_text.text = texto[0:48]
                    texto_pptx1_text.font.name = 'Helvetica'
                    texto_pptx1_text.font.size = Mm(10)

                    pdf.drawString(105*mm,25*mm, texto[48:len(texto)])

                    texto_pptx2 = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                    texto_pptx2_text_frame = texto_pptx2.text_frame
                    texto_pptx2_text_frame.clear()
                    texto_pptx2_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    texto_pptx2_text = texto_pptx2_text_frame.paragraphs[0].add_run()
                    texto_pptx2_text.text = texto[48:len(texto)]
                    texto_pptx2_text.font.name = 'Helvetica'
                    texto_pptx2_text.font.size = Mm(10)

                else:
                    pdf.drawString(105*mm,40*mm, capa['nome'])
                    
                    capa_nome = slide.shapes.add_textbox(Mm(102), Mm(170), Mm(200), Mm(10))
                    capa_nome_text_frame = capa_nome.text_frame
                    capa_nome_text_frame.clear()
                    capa_nome_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    capa_nome_text = capa_nome_text_frame.paragraphs[0].add_run()
                    capa_nome_text.text = capa['nome']
                    capa_nome_text.font.name = 'Helvetica'
                    capa_nome_text.font.size = Mm(10)

                    pdf.drawString(105*mm,25*mm, texto)

                    texto_pptx1 = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                    texto_pptx1_text_frame = texto_pptx1.text_frame
                    texto_pptx1_text_frame.clear()
                    texto_pptx1_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    texto_pptx1_text = texto_pptx1_text_frame.paragraphs[0].add_run()
                    texto_pptx1_text.text = texto
                    texto_pptx1_text.font.name = 'Helvetica'
                    texto_pptx1_text.font.size = Mm(10)

            else:
                pdf.drawString(105*mm,40*mm, capa['nome'])

                capa_nome = slide.shapes.add_textbox(Mm(102), Mm(170), Mm(200), Mm(10))
                capa_nome_text_frame = capa_nome.text_frame
                capa_nome_text_frame.clear()
                capa_nome_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                capa_nome_text = capa_nome_text_frame.paragraphs[0].add_run()
                capa_nome_text.text = capa['nome']
                capa_nome_text.font.name = 'Helvetica'
                capa_nome_text.font.size = Mm(10)
                if capa['cliente'] and not capa['pessoa']:
                    texto = f'{capa["cliente"]}'
                    pdf.drawString(105*mm,25*mm, texto)

                    texto_pptx1 = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                    texto_pptx1_text_frame = texto_pptx1.text_frame
                    texto_pptx1_text_frame.clear()
                    texto_pptx1_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    texto_pptx1_text = texto_pptx1_text_frame.paragraphs[0].add_run()
                    texto_pptx1_text.text = texto
                    texto_pptx1_text.font.name = 'Helvetica'
                    texto_pptx1_text.font.size = Mm(10)

                elif capa['pessoa'] and not capa['cliente']:
                    texto = f'A\C {capa["pessoa"]}'
                    pdf.drawString(105*mm,25*mm, texto)

                    texto_pptx1 = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                    texto_pptx1_text_frame = texto_pptx1.text_frame
                    texto_pptx1_text_frame.clear()
                    texto_pptx1_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    texto_pptx1_text = texto_pptx1_text_frame.paragraphs[0].add_run()
                    texto_pptx1_text.text = texto
                    texto_pptx1_text.font.name = 'Helvetica'
                    texto_pptx1_text.font.size = Mm(10)

    pdf.showPage()

    for i, linha in enumerate(linhas):
        with open(f'app/static/media/pdf_provider_images/temp_image{i}.png', 'wb') as nova_imagem:
            imagem = requests.get(linha[foto_column], stream=True)
            if not imagem.ok:
                return False, 'Link de imagem inexistente.'
            else:
                for dado in imagem.iter_content():
                    nova_imagem.write(dado)


        coordenadas = f'{linha[latitude_column]}  {linha[longitude_column]}'

        # Inicializando PPT
        slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
        
        # Template de fundo PDF
        pdf.drawImage('app/static/media/pdf_provider_images/template_pdf.jpg', 0, 0, 400*mm, 220*mm)

        # Template de fundo PPTX
        template_fundo = slide.shapes.add_picture('app/static/media/pdf_provider_images/template_pdf.jpg', Mm(0), Mm(0), height=Mm(220), width=Mm(400))

        # Endereço PDF
        pdf.setFont('Helvetica-Bold', 7*mm)
        pdf.setFillColor(colors.white)
        pdf.drawCentredString(200*mm,207*mm, linha[endereco_column])

        # Endereço PPTX
        endereco = slide.shapes.add_textbox(Mm(0), Mm(5), Mm(420), Mm(5))
        endereco_text_frame = endereco.text_frame
        endereco_text_frame.clear()
        endereco_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        endereco_text = endereco_text_frame.paragraphs[0].add_run()
        endereco_text.text = linha[endereco_column]
        endereco_text.font.name = 'Helvetica'
        endereco_text.font.bold = True
        endereco_text.font.size = Mm(7)
        endereco_text.font.color.rgb = RGBColor(255,255,255)

        # Imagem PDF
        pdf.drawImage(f'app/static/media/pdf_provider_images/temp_image{i}.png', 3*mm, 18.19*mm, 305*mm, 180*mm)

        #Imagem PPTX
        imagem = slide.shapes.add_picture(f'app/static/media/pdf_provider_images/temp_image{i}.png', Mm(3), Mm(21.81), height=Mm(180), width=Mm(310))

        # Coordenadas PDF
        pdf.setFont('Helvetica-Bold', 5*mm)
        pdf.setFillColor(colors.black)
        pdf.drawString(4*mm, 10*mm, 'Coordenadas:')
        
        pdf.setFont('Helvetica', 5*mm)
        pdf.drawString(41.5*mm, 10*mm, coordenadas)

        # Coordenadas PPTX
        coordenadas_pptx = slide.shapes.add_textbox(Mm(3), Mm(205), Mm(30), Mm(5))
        coordenadas_text_frame = coordenadas_pptx.text_frame
        coordenadas_text_frame.clear()
        coordenadas_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        coordenadas_text = coordenadas_text_frame.paragraphs[0].add_run()
        coordenadas_text.text = 'Coordenadas:'
        coordenadas_text.font.name = 'Helvetica'
        coordenadas_text.font.bold = True
        coordenadas_text.font.size = Mm(5)

        coordenadas_content = slide.shapes.add_textbox(Mm(40.5), Mm(205), Mm(50), Mm(6))
        coordenadas_content_text_frame = coordenadas_content.text_frame
        coordenadas_content_text_frame.clear()
        coordenadas_content_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        coordenadas_content_text = coordenadas_content_text_frame.paragraphs[0].add_run()
        coordenadas_content_text.text = coordenadas
        coordenadas_content_text.font.name = 'Helvetica'
        coordenadas_content_text.font.size = Mm(5)
        
        # Código PDF
        pdf.setFont('Helvetica-Bold', 5*mm)
        pdf.drawString(180*mm, 10*mm, 'Código:')
        
        pdf.setFont('Helvetica', 5*mm)
        pdf.drawString(202*mm, 10*mm, linha[codigo_column])

        # Código PPTX
        codigo = slide.shapes.add_textbox(Mm(180), Mm(205), Mm(30), Mm(5))
        codigo_text_frame = codigo.text_frame
        codigo_text_frame.clear()
        codigo_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        codigo_text = codigo_text_frame.paragraphs[0].add_run()
        codigo_text.text = 'Codigo:'
        codigo_text.font.name = 'Helvetica'
        codigo_text.font.bold = True
        codigo_text.font.size = Mm(5)

        codigo_content = slide.shapes.add_textbox(Mm(202), Mm(205), Mm(50), Mm(6))
        codigo_content_text_frame = codigo_content.text_frame
        codigo_content_text_frame.clear()
        codigo_content_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        codigo_content_text = codigo_content_text_frame.paragraphs[0].add_run()
        codigo_content_text.text = linha[codigo_column]
        codigo_content_text.font.name = 'Helvetica'
        codigo_content_text.font.size = Mm(5)

        eixo_y_pdf = 194
        eixo_y_pptx = 20

        for coluna in other_columns:
            # Outras Colunas PDF
            pdf.setFont('Helvetica-Bold', 5*mm)
            pdf.drawString(312*mm, eixo_y_pdf*mm, str(coluna))
            
            pdf.setFont('Helvetica', 5*mm)
            pdf.drawString(312*mm, (eixo_y_pdf - 7) * mm, str(linha[coluna]))
            eixo_y_pdf -= 20

            # Outras colunas PPTX
            outros = slide.shapes.add_textbox(Mm(315), Mm(eixo_y_pptx), Mm(30), Mm(6))
            outros_text_frame = outros.text_frame
            outros_text_frame.clear()
            outros_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            outros_text_frame.paragraphs[0].word_wrap = True
            outros_text = outros_text_frame.paragraphs[0].add_run()
            outros_text.text = str(coluna)
            outros_text.font.name = 'Helvetica'
            outros_text.font.bold = True
            outros_text.font.size = Mm(5)

            outros_content = slide.shapes.add_textbox(Mm(315), Mm(eixo_y_pptx + 7), Mm(30), Mm(6))
            outros_content_text_frame = outros_content.text_frame
            outros_content_text_frame.clear()
            outros_content_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            outros_content_text_frame.paragraphs[0].word_wrap = True
            outros_content_text = outros_content_text_frame.paragraphs[0].add_run()
            outros_content_text.text = str(linha[coluna])
            outros_content_text.font.name = 'Helvetica'
            outros_content_text.font.size = Mm(5)
            eixo_y_pptx += 20
        
        # Página PDF
        pdf.setFont('Helvetica-Bold', 8*mm)
        pdf.setFillColor(colors.white)
        pdf.drawCentredString(387*mm, 6.5*mm, str(pdf.getPageNumber()))

        # Página PPTX
        pagina = slide.shapes.add_textbox(Mm(380), Mm(203.5), Mm(10), Mm(10))
        pagina_text_frame = pagina.text_frame
        pagina_text_frame.clear()
        pagina_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        pagina_text = pagina_text_frame.paragraphs[0].add_run()
        pagina_text.text = str(pdf.getPageNumber())
        pagina_text.font.name = 'Helvetica'
        pagina_text.font.bold = True
        pagina_text.font.size = Mm(8)
        pagina_text.font.color.rgb = RGBColor(255,255,255)
        
        pdf.showPage()
    pdf.save()
    apresentacao.save(f'app/static/media/pptx/{image_id}.pptx')
    return True, 'Arquivos PDF e PPTX gerados com sucesso.'
    '''except UnidentifiedImageError as error:
        print(str(error))
        return False, 'A planilha contém imagem com formato inválido. Você deve fornecer um link de imagem .jpg, .jpeg ou .png. Telas de Google Maps ou Street View não são aceitas.'
    except Exception as error:
        print(str(error))
        return False, 'Erro no servidor. Tente novamente.'
        '''
    